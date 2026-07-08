from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from sem_archive.db.repository import Repository
from sem_archive.models import AppSettings, FolderRecord, SemCase
from sem_archive.utils.image_io import prepare_for_pptx
from sem_archive.utils.paths import iter_image_files, relative_label


@dataclass
class ImageItem:
    path: Path
    label: str
    alt_text: str
    row_key: str
    page_key: str


class PptExportService:
    def __init__(self, repo: Repository, settings: AppSettings) -> None:
        self.repo = repo
        self.settings = settings

    def collect_images(
        self,
        sem: SemCase,
        selected_folder_ids: list[int] | None = None,
        selected_relative_paths: list[str] | None = None,
    ) -> list[ImageItem]:
        assert sem.id is not None
        sem_root = Path(sem.local_path)
        folders = self.repo.list_folders(sem.id)
        folder_by_rel = {f.relative_path: f for f in folders}
        lots = ", ".join(l.lot_no for l in self.repo.list_lots(sem.id))
        extensions = self.settings.extension_set()

        target_dirs: list[Path] = []
        if selected_folder_ids:
            id_set = set(selected_folder_ids)
            for folder in folders:
                if folder.id in id_set:
                    target_dirs.append(sem_root / folder.relative_path)
        elif selected_relative_paths:
            for rel in selected_relative_paths:
                target_dirs.append(sem_root / rel if rel else sem_root)
        else:
            target_dirs.append(sem_root)

        # 重複除去しつつ配下画像を収集
        seen: set[Path] = set()
        items: list[ImageItem] = []
        for directory in target_dirs:
            for image_path in iter_image_files(directory, extensions):
                resolved = image_path.resolve()
                if resolved in seen:
                    continue
                seen.add(resolved)
                rel = relative_label(sem_root, image_path)
                parent_rel = image_path.parent.relative_to(sem_root).as_posix() if image_path.parent != sem_root else ""
                folder_meta = self._resolve_folder_meta(parent_rel, folder_by_rel)
                slot_id = self._resolve_slot(parent_rel, folder_by_rel)
                condition = folder_meta.condition if folder_meta and folder_meta.condition else sem.condition
                page_key, row_key = self._page_row_keys(parent_rel, folder_by_rel)
                alt = " / ".join(
                    part
                    for part in [
                        condition or "(条件なし)",
                        f"Lot:{lots or '-'}",
                        f"Slot:{slot_id or '-'}",
                        f"Folder:{image_path.parent.name}",
                        f"File:{image_path.name}",
                    ]
                )
                items.append(
                    ImageItem(
                        path=image_path,
                        label=rel,
                        alt_text=alt,
                        row_key=row_key,
                        page_key=page_key,
                    )
                )
        return items

    def export(
        self,
        output_path: Path,
        semis: list[SemCase],
        selection_map: dict[int, list[int]] | None = None,
    ) -> Path:
        """selection_map: sem_case_id -> folder_ids。空ならSEM全体。"""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        cache_dir = output_path.parent / "_tiff_cache"

        for sem in semis:
            assert sem.id is not None
            folder_ids = (selection_map or {}).get(sem.id)
            items = self.collect_images(sem, selected_folder_ids=folder_ids)
            if not items:
                continue
            pages = self._group_by_page(items)
            for page_key, page_items in pages:
                rows = self._group_by_row(page_items)
                self._render_page_groups(prs, sem.request_no, page_key, rows, cache_dir)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        return output_path

    def _render_page_groups(
        self,
        prs: Presentation,
        request_no: str,
        page_key: str,
        rows: list[tuple[str, list[ImageItem]]],
        cache_dir: Path,
    ) -> None:
        per_row = max(1, self.settings.images_per_row)
        # 1スライドに何行入るか（タイトル＋ラベル領域を考慮）
        max_rows_per_slide = 3
        slide = None
        row_on_slide = 0

        for row_key, row_items in rows:
            # 行内を per_row でチャンク
            for start in range(0, len(row_items), per_row):
                chunk = row_items[start : start + per_row]
                if slide is None or row_on_slide >= max_rows_per_slide:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    title = f"{request_no}"
                    if page_key:
                        title += f" / {page_key}"
                    self._add_title(slide, title)
                    row_on_slide = 0
                self._draw_row(slide, chunk, row_on_slide, per_row, cache_dir, row_key)
                row_on_slide += 1

    @staticmethod
    def _add_title(slide, text: str) -> None:
        box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.35))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.bold = True

    def _draw_row(
        self,
        slide,
        items: list[ImageItem],
        row_index: int,
        per_row: int,
        cache_dir: Path,
        row_key: str,
    ) -> None:
        margin_x = Inches(0.3)
        top0 = Inches(0.55)
        usable_width = Inches(13.333) - margin_x * 2
        cell_w = usable_width / per_row
        row_h = Inches(2.15)
        image_h = Inches(1.55)
        label_h = Inches(0.45)
        top = top0 + row_h * row_index

        if row_key and self.settings.export_row_mode == "subdir":
            # 行見出し
            box = slide.shapes.add_textbox(margin_x, top - Inches(0.18), usable_width, Inches(0.2))
            tf = box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = row_key
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(60, 60, 60)

        for i, item in enumerate(items):
            left = margin_x + cell_w * i
            embed_path = prepare_for_pptx(item.path, cache_dir)
            try:
                pic = slide.shapes.add_picture(
                    str(embed_path),
                    left + Inches(0.05),
                    top,
                    width=cell_w - Inches(0.1),
                    height=image_h,
                )
                # アスペクト比維持のためいったん入れてから調整
                self._fit_picture(pic, left + Inches(0.05), top, cell_w - Inches(0.1), image_h)
                self._set_alt_text(pic, item.alt_text)
            except Exception:
                # 読めない画像はスキップしてラベルだけ
                pass
            label_box = slide.shapes.add_textbox(
                left,
                top + image_h,
                cell_w,
                label_h,
            )
            tf = label_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item.label
            p.font.size = Pt(7)
            p.alignment = PP_ALIGN.CENTER

    @staticmethod
    def _fit_picture(pic, left, top, max_w, max_h) -> None:
        width = int(pic.width)
        height = int(pic.height)
        if width <= 0 or height <= 0:
            return
        max_w_i = int(max_w)
        max_h_i = int(max_h)
        scale = min(max_w_i / width, max_h_i / height)
        new_w = int(width * scale)
        new_h = int(height * scale)
        pic.width = new_w
        pic.height = new_h
        pic.left = int(left) + (max_w_i - new_w) // 2
        pic.top = int(top) + (max_h_i - new_h) // 2

    @staticmethod
    def _set_alt_text(picture, text: str) -> None:
        try:
            pic = picture._element
            nv_pr = pic.nvPicPr.cNvPr
            nv_pr.set("descr", text)
            # newer powerpoint also uses a:extLst; descr is enough for alt text
            if "{http://schemas.openxmlformats.org/drawingml/2006/main}descr" not in nv_pr.attrib:
                pass
            nv_pr.set(qn("a:descr") if False else "descr", text)
        except Exception:
            pass

    def _group_by_page(self, items: list[ImageItem]) -> list[tuple[str, list[ImageItem]]]:
        groups: dict[str, list[ImageItem]] = {}
        order: list[str] = []
        for item in items:
            key = item.page_key
            if key not in groups:
                groups[key] = []
                order.append(key)
            groups[key].append(item)
        return [(k, groups[k]) for k in order]

    def _group_by_row(self, items: list[ImageItem]) -> list[tuple[str, list[ImageItem]]]:
        if self.settings.export_row_mode == "none":
            return [("", items)]
        groups: dict[str, list[ImageItem]] = {}
        order: list[str] = []
        for item in items:
            key = item.row_key
            if key not in groups:
                groups[key] = []
                order.append(key)
            groups[key].append(item)
        return [(k, groups[k]) for k in order]

    def _page_row_keys(
        self, parent_rel: str, folder_by_rel: dict[str, FolderRecord]
    ) -> tuple[str, str]:
        parts = [p for p in parent_rel.split("/") if p]
        page_mode = self.settings.export_page_mode
        row_mode = self.settings.export_row_mode

        if page_mode == "flat":
            page_key = ""
        else:
            # slot モード: パス上で is_slot な最初のフォルダをページキーに
            page_key = ""
            for i, _name in enumerate(parts):
                rel = "/".join(parts[: i + 1])
                meta = folder_by_rel.get(rel)
                if meta and meta.is_slot:
                    page_key = meta.slot_id or meta.folder_name
                    break
            if not page_key and parts:
                # slotが見つからなければ第1階層
                page_key = parts[0]

        if row_mode == "none":
            row_key = ""
        else:
            # ページキーより下の部分を行キーに（例: C/M/E）
            row_key = ""
            if page_mode != "flat" and page_key:
                # page に使った階層の直後
                used = 0
                for i, name in enumerate(parts):
                    rel = "/".join(parts[: i + 1])
                    meta = folder_by_rel.get(rel)
                    if meta and (meta.slot_id == page_key or meta.folder_name == page_key):
                        used = i + 1
                        break
                rest = parts[used:]
                row_key = "/".join(rest) if rest else ""
            else:
                row_key = "/".join(parts[1:]) if len(parts) > 1 else (parts[0] if parts else "")
        return page_key, row_key

    @staticmethod
    def _resolve_folder_meta(
        parent_rel: str, folder_by_rel: dict[str, FolderRecord]
    ) -> FolderRecord | None:
        if parent_rel in folder_by_rel:
            return folder_by_rel[parent_rel]
        # 上位へ辿る
        rel = parent_rel
        while rel:
            if rel in folder_by_rel:
                return folder_by_rel[rel]
            if "/" not in rel:
                break
            rel = rel.rsplit("/", 1)[0]
        return None

    @staticmethod
    def _resolve_slot(parent_rel: str, folder_by_rel: dict[str, FolderRecord]) -> str | None:
        parts = [p for p in parent_rel.split("/") if p]
        for i in range(len(parts), 0, -1):
            rel = "/".join(parts[:i])
            meta = folder_by_rel.get(rel)
            if meta and meta.slot_id:
                return meta.slot_id
        return None
